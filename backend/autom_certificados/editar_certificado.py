from pptx import Presentation
import os


def substituir_texto_pptx(
    caminho: str,
    nomearquivo: str,
    substituicoes: dict,
    novonome: str = None
):
    """
    Substitui textos no PPTX conforme dicionário de substituições.
    """
    caminhoarquivo = os.path.join(caminho, f"{nomearquivo}.pptx")
    prs = Presentation(caminhoarquivo)

    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, 'text_frame') and shape.text_frame is not None:
                for paragraph in shape.text_frame.paragraphs:
                    for key, value in substituicoes.items():
                        if key in paragraph.text:
                            paragraph.text = paragraph.text.replace(key, value)

    # Salvar com novo nome ou sobrescrever
    if novonome:
        prs.save(os.path.join(caminho, f"{novonome}.pptx"))
    else:
        prs.save(caminhoarquivo)


def lerpptx(caminho: str = "", nomearquivo: str = ""):
    """
    Lê o conteúdo de um arquivo PPTX.
    """
    if caminho != "" and nomearquivo != "":
        caminhoarquivo = os.path.join(caminho, f"{nomearquivo}.pptx")
        prs = Presentation(caminhoarquivo)
        conteudo = ""
        for slide in prs.slides:
            for shape in slide.shapes:
                # Verifica se o shape tem o atributo text_frame
                if hasattr(shape, 'text_frame'):
                    text_frame = getattr(shape, 'text_frame', None)
                    if text_frame is not None:
                        for paragraph in text_frame.paragraphs:
                            conteudo += paragraph.text + "\n"
        return conteudo
    return ""


# Exemplo de uso:
substituicoes = {
    "XXXXXX": "João da Silva",
    "XXX.XXX.XXX-XX": "123.456.789-00",
    "XX de XXXXX de XXXX": "01 de Julho de 2025",
    # Adicione outros marcadores aqui
}

substituir_texto_pptx(
    caminho="./",
    nomearquivo="certificado",
    substituicoes=substituicoes,
    novonome="certificado_editado"
)

# Verificar o resultado
print("=== ARQUIVO ORIGINAL ===")
print(lerpptx(caminho="./", nomearquivo="certificado"))

print("\n=== ARQUIVO EDITADO ===")
print(lerpptx(caminho="./", nomearquivo="certificado_editado"))

